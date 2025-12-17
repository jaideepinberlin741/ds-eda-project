"""
Script to create a PowerPoint presentation for Larry Sanders housing analysis
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml import color

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
TITLE_COLOR = color.RGBColor(31, 78, 121)  # Dark blue
SUBTITLE_COLOR = color.RGBColor(68, 114, 196)  # Medium blue

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    subtitle_shape.text = subtitle
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = SUBTITLE_COLOR
    
    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    text_frame = content_shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    
    for item in content_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(12)
    
    return slide

def add_two_column_slide(prs, title, left_items, right_items):
    """Add a slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    # Remove default placeholder
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            sp = shape.element
            sp.getparent().remove(sp)
    
    # Add left text box
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for item in left_items:
        p = left_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(16)
        p.space_after = Pt(10)
    
    # Add right text box
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4.5), Inches(5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for item in right_items:
        p = right_frame.add_paragraph()
        p.level = 0
        p.text = item
        p.font.size = Pt(16)
        p.space_after = Pt(10)
    
    return slide

# Slide 1: Title
add_title_slide(
    prs,
    "Housing Analysis for Larry Sanders",
    "King County Housing Market - Data-Driven Recommendations"
)

# Slide 2: Client Profile
add_content_slide(
    prs,
    "Client Profile: Larry Sanders",
    [
        "Buyer seeking waterfront property",
        "Limited budget constraints",
        "Requires central location in Seattle",
        "Prefers isolated neighborhood (minimal children)",
        "Goal: Find the perfect match for his family's needs"
    ]
)

# Slide 3: Project Overview
add_content_slide(
    prs,
    "Project Overview",
    [
        "Analyzed King County housing sales data",
        "Processed 21,000+ house sales records",
        "Engineered custom metrics for client requirements",
        "Identified optimal properties matching criteria",
        "Generated actionable recommendations"
    ]
)

# Slide 4: Methodology - Data Preparation
add_content_slide(
    prs,
    "Data Preparation & Cleaning",
    [
        "Merged house sales and house details datasets",
        "Handled missing values (waterfront, view, renovation year)",
        "Created price per square foot metric",
        "Validated data quality and consistency",
        "Prepared geographic coordinates for location analysis"
    ]
)

# Slide 5: Feature Engineering
add_two_column_slide(
    prs,
    "Custom Metrics Developed",
    [
        "Population Density",
        "• Measures housing density in area",
        "• Based on proximity to other houses",
        "",
        "Centrality to Seattle",
        "• Distance from downtown Seattle",
        "• Higher score = more central location",
        "• Steep penalty for distance"
    ],
    [
        "Kid Density",
        "• Proportion of multi-bedroom homes nearby",
        "• Lower score = fewer families with children",
        "",
        "Larry Index",
        "• Custom scoring metric",
        "• Balances centrality vs. isolation",
        "• Formula: Centrality / (1 + Population + Kid Density)"
    ]
)

# Slide 6: Key Insights - Geographic
add_content_slide(
    prs,
    "Key Insight #1: Geographic Distribution",
    [
        "Waterfront properties are concentrated in specific areas",
        "Central Seattle locations show higher property values",
        "Optimal matches found in specific geographic clusters",
        "35 properties match all strict criteria simultaneously",
        "Geographic clustering suggests neighborhood-specific opportunities"
    ]
)

# Slide 7: Key Insights - Price Analysis
add_content_slide(
    prs,
    "Key Insight #2: Price Patterns",
    [
        "Waterfront properties command premium prices",
        "Centrality significantly impacts property values",
        "Price per square foot varies by location and features",
        "Budget constraints limit options but quality matches exist",
        "Value opportunities identified in specific neighborhoods"
    ]
)

# Slide 8: Key Insights - Neighborhood Characteristics
add_content_slide(
    prs,
    "Key Insight #3: Neighborhood Composition",
    [
        "Kid density varies significantly across neighborhoods",
        "Isolation from families with children is achievable",
        "Central locations can still offer privacy",
        "Waterfront + central + low kid density = rare combination",
        "Larry Index successfully identifies optimal properties"
    ]
)

# Slide 9: Findings Summary
add_content_slide(
    prs,
    "Analysis Results",
    [
        "✓ 35 properties match all criteria: waterfront, central, isolated",
        "✓ Properties identified using Larry Index ranking",
        "✓ Geographic visualization confirms optimal locations",
        "✓ Price ranges suitable for limited budget identified",
        "✓ Ready-to-review property list generated"
    ]
)

# Slide 10: Recommendation 1
add_content_slide(
    prs,
    "Recommendation #1: Focus on Top-Ranked Properties",
    [
        "Prioritize properties with highest Larry Index scores",
        "These properties offer best balance of all requirements",
        "Review top 10-15 properties first for efficiency",
        "Consider properties slightly above budget if they offer exceptional value",
        "Schedule viewings for top candidates immediately"
    ]
)

# Slide 11: Recommendation 2
add_content_slide(
    prs,
    "Recommendation #2: Geographic Strategy",
    [
        "Focus search on identified geographic clusters",
        "Properties in these areas meet multiple criteria simultaneously",
        "Consider expanding search radius slightly if needed",
        "Use geographic visualization to understand neighborhood context",
        "Leverage location insights for negotiation"
    ]
)

# Slide 12: Recommendation 3
add_content_slide(
    prs,
    "Recommendation #3: Timing & Action Plan",
    [
        "Act quickly - properties matching all criteria are limited",
        "Monitor market for new listings in identified areas",
        "Consider properties that meet 3 of 4 criteria as alternatives",
        "Use Larry Index to evaluate new listings as they appear",
        "Maintain flexibility on one criterion if exceptional property found"
    ]
)

# Slide 13: Key Statistics
add_content_slide(
    prs,
    "Key Statistics & Results",
    [
        "✓ 35 properties identified matching all criteria",
        "✓ Analysis based on 21,000+ house sales records",
        "✓ Properties span multiple price ranges for budget flexibility",
        "✓ All properties verified: waterfront + central + low kid density",
        "✓ Larry Index calculated for optimal ranking and prioritization"
    ]
)

# Slide 14: Next Steps
add_content_slide(
    prs,
    "Next Steps",
    [
        "Review detailed property list with house IDs and specifications",
        "Schedule property viewings for top-ranked candidates",
        "Validate on-site: waterfront access, neighborhood feel, kid density",
        "Use analysis insights to inform negotiation strategy",
        "Continue monitoring market for additional opportunities"
    ]
)

# Slide 15: Assumptions & Methodology Notes
add_content_slide(
    prs,
    "Methodology Notes & Assumptions",
    [
        "Kid density calculated as proportion of multi-bedroom homes in radius",
        "Centrality based on distance from Seattle downtown (47.6062°N, 122.3321°W)",
        "Waterfront status from property records (may need verification)",
        "Budget constraints applied based on price distribution analysis",
        "All metrics normalized to 0-1 scale for comparability"
    ]
)

# Slide 16: Thank You
add_title_slide(
    prs,
    "Thank You",
    "Questions & Discussion"
)

# Save presentation
output_path = 'Larry_Sanders_Housing_Analysis.pptx'
prs.save(output_path)
print(f"✅ Presentation created successfully: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
print(f"\n💡 To export as PDF:")
print(f"   Open the PPTX file in PowerPoint/Google Slides and export as PDF")

